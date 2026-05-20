"""vti-slide-pptx-exporter — render an HTML deck into a .pptx.

Three public entry points:

    render_slides_to_pngs(html_path, out_dir, scale=2, force=False) -> list[Path]
    build_pptx(image_paths, notes, out_pptx, deck_title=None)       -> Path
    export_deck(html_path, out_pptx, **kwargs)                      -> Path

Layout assumption: every slide is a `<section class="slide ...">` whose
CSS box is 1280×720 (the VTI canvas tokens `--sw` / `--sh`). The
exporter screenshots each element at that native size, multiplied by
`scale` for sharper text in the PPTX.

Notes are written into each slide's notes pane. See SKILL.md for the
expected `{idx: {script, hints}}` shape.
"""

from __future__ import annotations

import json
import re
import textwrap
from pathlib import Path
from typing import Any, Iterable

# python-pptx is part of the repo .venv — fail at import time if missing.
from pptx import Presentation
from pptx.util import Inches, Emu


# -- Canvas constants --------------------------------------------------------

# Matches the page-builder tokens (`--sw: 1280px; --sh: 720px;`). If the
# canvas tokens ever change, update both here and in tokens.css.
SLIDE_W_PX = 1280
SLIDE_H_PX = 720

# 16:9 PPTX canvas. python-pptx uses EMU internally; Inches() converts.
# 13.333" × 7.5" = the PowerPoint 16:9 widescreen default.
PPTX_W_IN = 13.333
PPTX_H_IN = 7.5


# -- Slide enumeration -------------------------------------------------------

# Match top-level <div class="slide ..."> (the page-builder root for
# each slide). `class` must START with `slide` so we don't accidentally
# pick up `vti-slide-*` component wrappers.
_SLIDE_RE = re.compile(
    r'<div\b[^>]*\bclass="slide(?:\s[^"]*)?"[^>]*>',
    re.IGNORECASE,
)

# CSS selector used by Playwright. `div.slide` matches the class
# `slide` exactly — not `vti-slide-*` — so it picks only the slide
# roots even if the markup ever adds more `.vti-slide-*` descendants.
_SLIDE_SELECTOR = "div.slide"


def count_slides(html: str) -> int:
    return len(_SLIDE_RE.findall(html))


def read_deck_title(html: str) -> str:
    m = re.search(r"<title>(.*?)</title>", html, re.IGNORECASE | re.DOTALL)
    return (m.group(1).strip() if m else "VTI Slide Deck")


# -- Rendering ---------------------------------------------------------------

_PLAYWRIGHT_MISSING_MSG = textwrap.dedent(
    """
    Playwright is not installed in this Python environment. The exporter
    needs a headless Chromium to screenshot each slide.

    Install once:
        pip install playwright
        python -m playwright install chromium

    Then re-run the export driver.
    """
).strip()


def _ensure_playwright():
    try:
        from playwright.sync_api import sync_playwright  # noqa: F401
    except ImportError as exc:
        raise RuntimeError(_PLAYWRIGHT_MISSING_MSG) from exc


def render_slides_to_pngs(
    html_path: str | Path,
    out_dir: str | Path,
    scale: int = 2,
    force: bool = False,
) -> list[Path]:
    """Screenshot each `.slide` section in the deck HTML.

    Returns the list of output PNG paths in deck order. Files are
    named `slide_01.png`, `slide_02.png`, …

    Caching: if every expected output already exists AND is newer than
    the deck HTML, the render is skipped. Pass `force=True` to override.
    """
    html_path = Path(html_path).resolve()
    out_dir = Path(out_dir).resolve()
    if not html_path.exists():
        raise FileNotFoundError(f"deck HTML not found: {html_path}")

    out_dir.mkdir(parents=True, exist_ok=True)

    html = html_path.read_text(encoding="utf-8")
    n_slides = count_slides(html)
    if n_slides == 0:
        raise RuntimeError(
            f"no <section class='slide ...'> found in {html_path}. "
            "Is this a composed deck?"
        )

    expected = [out_dir / f"slide_{i:02d}.png" for i in range(1, n_slides + 1)]

    if not force and _cache_is_fresh(expected, html_path):
        print(f"  ✓ pptx-images cache fresh ({n_slides} slides) — skipping render")
        return expected

    # Wipe stale PNGs so a shrunken deck doesn't leave orphans.
    for stale in out_dir.glob("slide_*.png"):
        stale.unlink()

    _ensure_playwright()
    from playwright.sync_api import sync_playwright

    print(f"  → rendering {n_slides} slides at {SLIDE_W_PX}×{SLIDE_H_PX} × {scale}…")

    with sync_playwright() as p:
        browser = p.chromium.launch()
        try:
            ctx = browser.new_context(
                viewport={"width": SLIDE_W_PX, "height": SLIDE_H_PX},
                device_scale_factor=scale,
            )
            page = ctx.new_page()
            page.goto(html_path.as_uri(), wait_until="networkidle")
            # Wait for fonts to settle so Plus Jakarta Sans is not
            # mid-swap when we screenshot.
            page.evaluate("() => document.fonts && document.fonts.ready")

            slides = page.locator(_SLIDE_SELECTOR)
            count = slides.count()
            if count != n_slides:
                # Defensive: regex count vs DOM count should agree; if
                # not, trust the DOM.
                print(f"  ! regex saw {n_slides} slides, DOM has {count} — using DOM count")
                expected = [out_dir / f"slide_{i:02d}.png" for i in range(1, count + 1)]

            for i in range(count):
                target = expected[i]
                slides.nth(i).screenshot(path=str(target))
                print(f"  ✓ slide_{i+1:02d} → {target.name}")
        finally:
            browser.close()

    return expected


def _cache_is_fresh(expected: Iterable[Path], html_path: Path) -> bool:
    expected = list(expected)
    if not expected:
        return False
    if not all(p.exists() for p in expected):
        return False
    src_mtime = html_path.stat().st_mtime
    return all(p.stat().st_mtime >= src_mtime for p in expected)


# -- PPTX assembly -----------------------------------------------------------

NOTES_TEMPLATE = textwrap.dedent(
    """
    ═══ TRÌNH BÀY ═══
    {script}

    ═══ GỢI Ý ═══
    {hints}
    """
).strip()


def _format_notes(entry: dict | None) -> str:
    if not entry:
        return ""
    script = (entry.get("script") or "").strip() or "(chưa có script)"
    hints = (entry.get("hints") or "").strip() or "(chưa có gợi ý)"
    return NOTES_TEMPLATE.format(script=script, hints=hints)


def _normalize_notes(notes: Any) -> dict[int, dict]:
    """Accept either {idx: {...}} dict or list-of-records → dict by 1-based index."""
    if notes is None:
        return {}
    if isinstance(notes, dict):
        return {int(k): v for k, v in notes.items()}
    if isinstance(notes, list):
        out: dict[int, dict] = {}
        for i, item in enumerate(notes, start=1):
            if not isinstance(item, dict):
                continue
            idx = int(item.get("slide_index") or i)
            out[idx] = item
        return out
    raise TypeError(f"notes must be dict or list, got {type(notes).__name__}")


def build_pptx(
    image_paths: list[Path],
    notes: Any,
    out_pptx: str | Path,
    deck_title: str | None = None,
) -> Path:
    """Assemble a 16:9 .pptx with one image-fill slide per PNG and
    notes pasted into each slide's notes pane.
    """
    out_pptx = Path(out_pptx).resolve()
    out_pptx.parent.mkdir(parents=True, exist_ok=True)

    if not image_paths:
        raise ValueError("image_paths is empty — nothing to assemble")

    notes_by_idx = _normalize_notes(notes)

    prs = Presentation()
    prs.slide_width = Inches(PPTX_W_IN)
    prs.slide_height = Inches(PPTX_H_IN)

    blank_layout = prs.slide_layouts[6]  # "Blank" — no placeholders

    for i, img in enumerate(image_paths, start=1):
        slide = prs.slides.add_slide(blank_layout)

        # Full-bleed image at (0,0). python-pptx sizes the picture to
        # the explicit width/height we pass; aspect-ratio of the source
        # PNG (1280×720) matches the slide (13.333"×7.5") so there is
        # no visible distortion.
        slide.shapes.add_picture(
            str(img),
            left=Emu(0),
            top=Emu(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )

        notes_text = _format_notes(notes_by_idx.get(i))
        if notes_text:
            slide.notes_slide.notes_text_frame.text = notes_text

    if deck_title:
        prs.core_properties.title = deck_title

    prs.save(str(out_pptx))
    return out_pptx


# -- Convenience entry point -------------------------------------------------

def export_deck(
    html_path: str | Path,
    out_pptx: str | Path | None = None,
    images_dir: str | Path | None = None,
    notes_path: str | Path | None = None,
    phase_2_path: str | Path | None = None,
    phase_3_path: str | Path | None = None,
    deck_title: str | None = None,
    scale: int = 2,
    force_render: bool = False,
    force_notes: bool = False,
) -> Path:
    """One-call export: render PNGs + seed/load notes + build PPTX.

    Steps:
      1. Render slides to PNGs (cached by HTML mtime).
      2. If `notes_path` does not exist (or `force_notes=True`) and the
         phase JSONs are reachable, seed a baseline `pptx-notes.json`
         from the outline + block content.
      3. Read `notes_path` if present; build the PPTX with those notes.

    `out_pptx` defaults to `<html_path with .pptx>` so the PPTX inherits
    the deck filename convention (see vti-slide-creator/deck_filename.py).

    Returns the path to the written .pptx.
    """
    html_path = Path(html_path).resolve()
    if out_pptx is None:
        out_pptx = html_path.with_suffix(".pptx")
    out_pptx = Path(out_pptx).resolve()
    images_dir = Path(images_dir).resolve() if images_dir else out_pptx.parent / "pptx-images"
    notes_path = Path(notes_path).resolve() if notes_path else out_pptx.parent / "pptx-notes.json"

    # 1. Render
    image_paths = render_slides_to_pngs(
        html_path=html_path,
        out_dir=images_dir,
        scale=scale,
        force=force_render,
    )

    # 2. Seed notes if needed
    if force_notes or not notes_path.exists():
        if phase_2_path or phase_3_path:
            try:
                from notes_builder import derive_notes_from_phases

                seed = derive_notes_from_phases(
                    phase_2_path=phase_2_path,
                    phase_3_path=phase_3_path,
                    n_slides=len(image_paths),
                )
                notes_path.parent.mkdir(parents=True, exist_ok=True)
                notes_path.write_text(
                    json.dumps(seed, ensure_ascii=False, indent=2),
                    encoding="utf-8",
                )
                action = "overwrote" if force_notes else "seeded"
                print(f"  ✓ {action} {notes_path.name} from phase outputs")
            except Exception as exc:
                print(f"  ! could not seed notes from phases: {exc}")
        else:
            print(f"  ! no phase paths supplied — {notes_path.name} not seeded")

    # 3. Load notes (whatever is on disk) and build
    notes: Any = None
    if notes_path.exists():
        try:
            notes = json.loads(notes_path.read_text(encoding="utf-8"))
        except Exception as exc:
            print(f"  ! failed to parse {notes_path.name} — building without notes ({exc})")

    if not deck_title:
        try:
            deck_title = read_deck_title(html_path.read_text(encoding="utf-8"))
        except Exception:
            deck_title = None

    written = build_pptx(
        image_paths=image_paths,
        notes=notes,
        out_pptx=out_pptx,
        deck_title=deck_title,
    )
    print(f"  ✓ wrote {written.relative_to(written.parents[1]) if len(written.parents) > 1 else written}")
    return written


# -- Allow `python -m exporter <html> <pptx>` style invocation ---------------

if __name__ == "__main__":
    import argparse

    ap = argparse.ArgumentParser(description="Export an HTML deck to PPTX with speaker notes.")
    ap.add_argument("html", help="Path to deck-composed.html")
    ap.add_argument("pptx", help="Output .pptx path")
    ap.add_argument("--images-dir", default=None)
    ap.add_argument("--notes", default=None, help="Path to pptx-notes.json")
    ap.add_argument("--phase-2", default=None)
    ap.add_argument("--phase-3", default=None)
    ap.add_argument("--scale", type=int, default=2)
    ap.add_argument("--force-render", action="store_true")
    ap.add_argument("--force-notes", action="store_true")
    args = ap.parse_args()

    export_deck(
        html_path=args.html,
        out_pptx=args.pptx,
        images_dir=args.images_dir,
        notes_path=args.notes,
        phase_2_path=args.phase_2,
        phase_3_path=args.phase_3,
        scale=args.scale,
        force_render=args.force_render,
        force_notes=args.force_notes,
    )
